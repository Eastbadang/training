# python-pptx 패키지 설치 : pip install python-pptx

from pptx import Presentation # 라이브러리
from pptx.util import Inches # 사진, 표등을 그리기 위해

# 제목 슬라이드 추가

prs = Presentation() # 파워포인트 객체 선언

title_slide_layout = prs.slide_layouts[0] # 0 : 제목슬라이드에 해당
slide = prs.slides.add_slide(title_slide_layout) # 제목 슬라이드를 파워포인트 객체에 추가

# 제목 - 제목에 값넣기
title = slide.shapes.title  # 제목
title.text = "Hello, World!" # 제목에 값 넣기

# 부제목
subtitle = slide.placeholders[1] # 제목상자는 placeholders[0], 부제목상자는 [1]
subtitle.text = "python-pptx was here!"

# 저장
prs.save('test1.pptx')

# 제목 + 내용 슬라이드 추가

bullet_slide_layout = prs.slide_layouts[1] # 1 : 제목 + 내용 슬라이드

slide = prs.slides.add_slide(bullet_slide_layout) # 기존에 있던 슬라이드에 추가
shapes = slide.shapes

# 제목
title_shape = shapes.title
title_shape.text = 'Adding a Bullet Slide'

# 내용
body_shape = shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

# 단락 추가
p = tf.add_paragraph()
p.text = 'Use _TextFrame.text for first bullet'
p.level = 1  # 1 : 들여쓰기 레벨

# 단락 추가
p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2  # 2 : 들여쓰기 레벨

prs.save('test2.pptx')

# 제목 + 내용 슬라이드 추가

bullet_slide_layout = prs.slide_layouts[1] # 1 : 제목 + 내용 슬라이드

slide = prs.slides.add_slide(bullet_slide_layout) # 기존에 있던 슬라이드에 추가
shapes = slide.shapes

# 제목
title_shape = shapes.title
title_shape.text = 'Adding a Bullet Slide'

# 내용
body_shape = shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

# 단락 추가
p = tf.add_paragraph()
p.text = 'Use _TextFrame.text for first bullet'
p.level = 1  # 1 : 들여쓰기 레벨

# 단락 추가
p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2  # 2 : 들여쓰기 레벨

prs.save('test3.pptx')